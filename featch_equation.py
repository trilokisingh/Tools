import docx
from pptx import Presentation
from pptx.util import Inches

def extract_equations_from_docx(docx_file):
    equations = []
    document = docx.Document(docx_file)

    for paragraph in document.paragraphs:
        equations.append(paragraph.text)

    return equations

def create_pptx_with_equations(equations, pptx_file):
    presentation = Presentation()

    for equation in equations:
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)

        left = top = width = height = Inches(1)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = equation

    presentation.save(pptx_file)
    print(f"Presentation saved as {pptx_file}")

# here we are giving input 
equations = extract_equations_from_docx("Sample_DOCX.docx")
create_pptx_with_equations(equations, "output.pptx")
